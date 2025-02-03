from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import json
import os
from datetime import datetime

def json_to_ppt(json_data: dict, filename: str = None) -> str:
    """
    Converts JSON pitch deck data to a PowerPoint presentation
    Returns path to created PPTX file
    """
    try:
        # Create default filename if none provided
        if filename is None:
            prompt = json_data.get('prompt', 'untitled')
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"pitch_deck_{prompt.replace(' ', '_')}_{timestamp}.pptx"

        prs = Presentation()
        
        title_slide = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Pitch Deck: " + json_data.get("prompt", "")
        title_paragraph = title.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(44)
        
        subtitle.text = "Generated on " + json_data.get("timestamp", "")
        subtitle_paragraph = subtitle.text_frame.paragraphs[0]
        subtitle_paragraph.font.size = Pt(24)

        for slide_data in json_data.get("slides", []):
            slide_layout = prs.slide_layouts[1]
            new_slide = prs.slides.add_slide(slide_layout)
            
            title_shape = new_slide.shapes.title
            title_shape.text = slide_data.get("title", "")
            title_frame = title_shape.text_frame
            title_frame.paragraphs[0].font.size = Pt(40)
            
            content_shape = new_slide.placeholders[1]
            content_frame = content_shape.text_frame
            content_text = slide_data.get("content", "")
            
            content_frame.clear()
            
            lines = content_text.split("\n")
            for line in lines:
                p = content_frame.add_paragraph()
                if line.startswith("-"):
                    p.text = line[1:].strip()  # Remove the dash
                    p.level = 0  # This makes it a bullet point
                else:
                    p.text = line
                
                p.font.size = Pt(24)  # Smaller font for content
                p.space_before = Pt(12)  # Add some spacing between points
                p.space_after = Pt(6)

        current_dir = os.path.dirname(os.path.abspath(__file__))
        output_dir = os.path.join(current_dir, '..', 'outputs')
        os.makedirs(output_dir, exist_ok=True)
        
        # Save presentation
        output_path = os.path.join(output_dir, filename)
        prs.save(output_path)
        
        return output_path
        
    except Exception as e:
        print(f"Error generating PowerPoint: {str(e)}")
        raise
